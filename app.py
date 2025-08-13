import os
from flask import Flask, render_template, request, jsonify, send_from_directory, flash, redirect, url_for
from flask_login import LoginManager, current_user, login_required
from werkzeug.utils import secure_filename
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
import string
import filetype
from pdfminer.high_level import extract_text
from pptx import Presentation
import docx
import csv
import logging
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import spacy
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from sentence_transformers import SentenceTransformer
import wikipediaapi
from collections import defaultdict

# Initialize NLTK
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)

# Initialize Flask
app = Flask(__name__, template_folder='../frontend/templates', static_folder='../frontend/static')
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'fallback-secret-key')
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///db.sqlite'
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'frontend', 'static', 'uploads')
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload size

# Initialize extensions
from database import db
db.init_app(app)

# Initialize Login Manager
login_manager = LoginManager()
login_manager.login_view = 'auth.login'
login_manager.init_app(app)

from database import User

@login_manager.user_loader
def load_user(user_id):
    return db.session.get(User, int(user_id))

# Register Blueprints
from auth import auth as auth_blueprint
app.register_blueprint(auth_blueprint)

app.logger.setLevel(logging.DEBUG)

# Load NLP models
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    raise RuntimeError("Please download the English language model first:\n"
                      "python -m spacy download en_core_web_sm")

class AdvancedChatbot:
    def __init__(self):
        # Initialize models
        self.sentence_model = SentenceTransformer('all-MiniLM-L6-v2')
        self.wiki = wikipediaapi.Wikipedia('en')
        self.tfidf = TfidfVectorizer()
        
        # Enhanced knowledge base
        self.knowledge_base = [
            {
                "patterns": ["hello", "hi", "hey"],
                "response": "Hello! How can I assist you today?",
                "category": "greeting"
            },
            {
                "patterns": ["help", "support"],
                "response": "I can help with document analysis and answer questions. Try asking about file uploads or general knowledge.",
                "category": "help"
            },
            {
                "patterns": ["password", "reset", "login"],
                "response": "You can reset your password at /reset-password or contact support@example.com",
                "category": "account"
            },
            {
                "patterns": ["file", "upload", "document"],
                "response": "You can upload PDF, DOCX, PPTX, TXT, and CSV files for analysis.",
                "category": "features"
            }
        ]
        
        # Precompute embeddings
        self.knowledge_embeddings = []
        all_patterns = []
        for item in self.knowledge_base:
            patterns = " ".join(item["patterns"])
            all_patterns.append(patterns)
            self.knowledge_embeddings.append(self.sentence_model.encode(patterns))
        
        # Fit TF-IDF
        self.tfidf.fit(all_patterns)
        self.similarity_threshold = 0.7
    
    def get_wikipedia_answer(self, query):
        try:
            page = self.wiki.page(query)
            if page.exists():
                summary = page.summary[:500]
                if len(page.summary) > 500:
                    summary += "... (see more at Wikipedia)"
                return summary
            return None
        except Exception as e:
            app.logger.error(f"Wikipedia API error: {str(e)}")
            return None
    
    def respond(self, message):
        try:
            message = message.lower().strip("?,.!")
            
            # Check for exact matches first
            for item in self.knowledge_base:
                if any(pattern in message for pattern in item["patterns"]):
                    return item["response"]
            
            # Semantic similarity matching
            query_embedding = self.sentence_model.encode(message)
            similarities = cosine_similarity([query_embedding], self.knowledge_embeddings)[0]
            best_match_idx = np.argmax(similarities)
            
            if similarities[best_match_idx] > self.similarity_threshold:
                return self.knowledge_base[best_match_idx]["response"]
            
            # Fallback to Wikipedia for general knowledge
            wiki_answer = self.get_wikipedia_answer(message)
            if wiki_answer:
                return f"Here's what I found:\n\n{wiki_answer}\n\nFor specific help with our services, try asking about document analysis or account support."
            
            # Final fallback
            return "I'm not sure I understand. Try asking about:\n- File uploads\n- Password help\n- General knowledge questions"
        except Exception as e:
            app.logger.error(f"Chatbot error: {str(e)}")
            return "Sorry, I'm having trouble responding right now. Please try again later."

# Initialize chatbot instance
chatbot_engine = AdvancedChatbot()

# File Processing Functions
def allowed_file(filename):
    ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'pptx', 'csv'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(filepath):
    try:
        kind = filetype.guess(filepath)
        if kind is None:
            ext = os.path.splitext(filepath)[1].lower()
            if ext == '.txt':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            return "Unsupported file type"
        
        if kind.mime == 'text/plain':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif kind.mime == 'application/pdf':
            return extract_text(filepath)
        elif kind.mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc = docx.Document(filepath)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif kind.mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            prs = Presentation(filepath)
            return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif kind.mime == 'text/csv':
            with open(filepath, 'r', encoding='utf-8') as f:
                return '\n'.join([','.join(row) for row in csv.reader(f)])
        else:
            return "Unsupported file type"
    except Exception as e:
        app.logger.error(f"Error processing {filepath}: {str(e)}")
        return f"Error processing file: {str(e)}"

def analyze_text(text):
    try:
        # Process with spaCy
        doc = nlp(text)
        
        # Extract keywords (noun phrases)
        keywords = [chunk.text for chunk in doc.noun_chunks][:10]
        
        # Summarization with sumy
        parser = PlaintextParser.from_string(text, Tokenizer("english"))
        summarizer = LsaSummarizer()
        summary = " ".join([str(sentence) for sentence in summarizer(parser.document, 3)])
        
        # Statistics
        tokens = [token.text for token in doc if not token.is_punct]
        
        return {
            'keywords': keywords,
            'summary': summary,
            'word_count': len(tokens),
            'unique_words': len(set(tokens))
        }
    except Exception as e:
        app.logger.error(f"Text analysis error: {str(e)}")
        return {
            'keywords': [],
            'summary': "Could not generate summary",
            'word_count': 0,
            'unique_words': 0
        }

def get_chatbot_response(message):
    return chatbot_engine.respond(message)

# Routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/about')
def about():
    return render_template('about.html')

@app.route('/contact')
def contact():
    return render_template('contact.html')

@app.route('/analytics')
@login_required
def analytics():
    return render_template('analytics.html')

@app.route('/file_upload', methods=['GET', 'POST'])
@login_required
def file_upload():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file part')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No selected file')
            return redirect(request.url)
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            
            try:
                os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
                file.save(filepath)
                
                if not os.path.exists(filepath):
                    flash('File failed to save')
                    return redirect(request.url)
                
                text = extract_text_from_file(filepath)
                if text == "Unsupported file type":
                    flash('Unsupported file type')
                    return redirect(request.url)
                    
                analysis = analyze_text(text)
                return render_template('file_upload.html', analysis=analysis, filename=filename)
                
            except Exception as e:
                flash(f'Error processing file: {str(e)}')
                app.logger.error(f"File processing error: {str(e)}")
                return redirect(request.url)
    
    return render_template('file_upload.html')

@app.route('/chatbot', methods=['GET', 'POST'])
@login_required
def chatbot():
    if request.method == 'POST':
        data = request.get_json()
        user_message = data.get('message', '')
        response = get_chatbot_response(user_message)
        return jsonify({'response': response})
    
    return render_template('chatbot.html')

@app.route('/uploads/<filename>')
@login_required
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/test')
def test_endpoints():
    """Test all functionality"""
    # Test chatbot
    tests = [
        ("hello", "greeting"),
        ("help with password", "account"),
        ("what files can I upload?", "features"),
        ("tell me about AI", "wikipedia")
    ]
    
    results = []
    for question, expected in tests:
        response = get_chatbot_response(question)
        results.append({
            "question": question,
            "response": response,
            "expected": expected,
            "status": "success" if expected in response.lower() else "warning"
        })
    
    # Test file processing
    test_file = os.path.join(app.root_path, 'test.txt')
    try:
        with open(test_file, 'w') as f:
            f.write("This is a test file for NLP processing.")
        file_content = extract_text_from_file(test_file)
        analysis = analyze_text(file_content)
        file_status = "success"
    except Exception as e:
        file_status = f"error: {str(e)}"
    
    return jsonify({
        "chatbot_tests": results,
        "file_processing": file_status,
        "models_loaded": True
    })

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True)